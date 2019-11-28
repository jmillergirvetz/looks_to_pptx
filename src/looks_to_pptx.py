# import modules
import sys
import os
import io

from looker_sdk import client, models, error
from pptx import Presentation
from pptx.util import Cm


# client calls will now automatically authenticate using the
# api3credentials specified in 'looker.ini' which is then added to .gitginore
sdk = client.setup("../looker.ini")


def get_pres_temp_layout(pptx_tmp_path):
    """Create a powerpoint presentation object from a pptx template"""
    try:
        prs = Presentation(pptx_tmp_path)
    except:
        print(f'Cannot locate \'{pptx_tmp_path}\'. Please check the path to the desired .pptx template.')
    
    for idx, layout in enumerate(prs.slide_layouts):
        print(idx, layout.name)
    return prs


def cnt_slides(prs):
    """Returns the total amount of slides, used to locate the ending slide from where the additional slides will be added"""
    cnt = 0
    for slide in prs.slides:
        cnt += 1
    return cnt


def looks_to_pptx(looks, prs, slide_layout, powerpoint_template, target_space):
    """Creates slides at the end of the powerpoint presentation and adds Looks from the specified target space"""
    ending_slide = cnt_slides(prs)
    
    for idx, look in enumerate(looks):
        print(idx, look.id, look.title, ' - Look Added')

        # looker's python sdk returns the requested png as an image byte stream for performant rendering results
        look_img_byte_string = sdk.run_look(look.id, 'png')
        # handles the look's image byte string and converts it into a temporary in-memory file
        tmpFile = io.BytesIO(look_img_byte_string)

        # adds slide to the end of the presentation with the specified slide layout - this should be a "Title Only" layout
        prs.slides.add_slide(prs.slide_layouts[slide_layout])

        # gets the shapes object from the most recent slide and adds the Look as a png with the specific spacing [Cm(5.4), Cm(4.95), width=Cm(23)]
        if ending_slide == 0:
            shapes = prs.slides[idx].shapes
            shapes.add_picture(tmpFile, Cm(2.4), Cm(2.95), width=Cm(20))
        else:
            shapes = prs.slides[ending_slide+idx].shapes
            shapes.add_picture(tmpFile, Cm(2.4), Cm(2.95), width=Cm(20))
    
    # writes powerpoint with the original powerpoint template file name and includes a tag for looker generated and the target space 
    base = os.path.basename(powerpoint_template)
    filename = os.path.splitext(base)[0]
    
    prs.save('./' + filename + f'(looker_generated_{target_space}).pptx')


def main():

    # shows you as authenticated
    looker_api_user = sdk.me()
    print(looker_api_user)
    
    try:
        target_space =  int(sys.argv[1]) if len(sys.argv) >= 1 else None
        template_slide = int(sys.argv[3]) if len(sys.argv) >= 3 else None
    except:
        print(sys.argv[1] + ' and/or ' + sys.argv[3] + ' is not a whole integer. Please enter the target space\'s id and the template slide as whole integers.')
    
    powerpoint_template = sys.argv[2] if len(sys.argv) >= 2 else None

    if not target_space or not template_slide or not powerpoint_template:
        print(f"Please provide: <targetSpace> <pptxTemplate> <templateSlide>")
        pass


    looks = sdk.space_looks(target_space)

    prs = get_pres_temp_layout(powerpoint_template)

    looks_to_pptx(looks, prs, template_slide, powerpoint_template, target_space)


main()